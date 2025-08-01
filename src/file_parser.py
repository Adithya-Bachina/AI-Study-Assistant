# src/file_parser.py

import PyPDF2
import docx
from pptx import Presentation
from io import BytesIO

def extract_text_from_pdf(file_stream):
    """Extracts text from a PDF file stream."""
    pdf_reader = PyPDF2.PdfReader(file_stream)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_stream):
    """Extracts text from a DOCX file stream."""
    document = docx.Document(file_stream)
    text = ""
    for paragraph in document.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(file_stream):
    """Extracts text from a PPTX file stream."""
    prs = Presentation(file_stream)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                text += shape.text_frame.text + "\n"
    return text